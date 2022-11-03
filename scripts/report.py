#!/usr/bin/env python3
import click
import pptx
from pptx.enum.shapes import MSO_SHAPE
from datetime import date
import os
import pandas as pd

NO_DATA_CHAR = "NA"
TITLE = "SARS-CoV-2 Recombinants"
FONT_SIZE_PARAGRAPH = 20

RECOMBINANT_STATUS = [
    "designated",
    "proposed",
    "unpublished",
]

""" Ref for slide types:
0 ->  title and subtitle
1 ->  title and content
2 ->  section header
3 ->  two content
4 ->  Comparison
5 ->  Title only
6 ->  Blank
7 ->  Content with caption
8 ->  Pic with caption
"""


@click.command()
@click.option("--linelist", help="Linelist TSV", required=True)
@click.option("--plot-dir", help="Plotting directory", required=True)
@click.option("--output", help="Output PPTX path", required=True)
@click.option(
    "--geo", help="Geography column to summarize", required=False, default="country"
)
@click.option(
    "--singletons",
    help="Whether singletons were included in plots",
    is_flag=True,
)
@click.option(
    "--template",
    help="Powerpoint template",
    required=False,
    default="resources/template.pptx",
)
def main(
    linelist,
    plot_dir,
    template,
    geo,
    output,
    singletons,
):
    """Create a report of powerpoint slides"""

    # Parse build name from plot_dir path
    build = os.path.basename(os.path.dirname(plot_dir))
    outdir = os.path.dirname(output)
    if not os.path.exists(outdir):
        os.mkdir(outdir)
    subtitle = "{}\n{}".format(build.title(), date.today())

    # Import dataframes
    linelist_df = pd.read_csv(linelist, sep="\t")

    plot_suffix = ".png"
    df_suffix = ".tsv"
    labels = [
        f.replace(plot_suffix, "")
        for f in os.listdir(plot_dir)
        if f.endswith(plot_suffix)
    ]

    plot_dict = {}
    for label in labels:
        plot_dict[label] = {}
        plot_dict[label]["plot_path"] = os.path.join(plot_dir, label + plot_suffix)
        plot_dict[label]["df_path"] = os.path.join(plot_dir, label + df_suffix)
        plot_dict[label]["df"] = pd.read_csv(plot_dict[label]["df_path"], sep="\t")

        # Breakpoints df isn't over time, but by lineage
        if "epiweek" in plot_dict[label]["df"].columns:
            plot_dict[label]["df"].index = plot_dict[label]["df"]["epiweek"]

        # Largest is special, as it takes the form largest_<lineage>.*
        if label.startswith("largest_"):

            largest_lineage = "_".join(label.split("_")[1:])
            # Replace the _DELIM_ character we added for saving files
            largest_lineage = largest_lineage.replace("_DELIM_", "/")

            plot_dict["largest"] = plot_dict[label]
            plot_dict["largest"]["lineage"] = largest_lineage

            del plot_dict[label]

    # ---------------------------------------------------------------------
    # Presentation
    # ---------------------------------------------------------------------
    presentation = pptx.Presentation(template)

    # ---------------------------------------------------------------------
    # Title Slide
    title_slide_layout = presentation.slide_layouts[0]
    slide = presentation.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = TITLE
    slide.placeholders[1].text = subtitle

    # ---------------------------------------------------------------------
    # General Summary

    # Slide setup
    plot_path = plot_dict["lineage"]["plot_path"]
    graph_slide_layout = presentation.slide_layouts[8]
    slide = presentation.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title
    title.text_frame.text = "Status"
    title.text_frame.paragraphs[0].font.bold = True

    chart_placeholder = slide.placeholders[1]
    chart_placeholder.insert_picture(plot_path)
    body = slide.placeholders[2]

    # Stats

    lineages_df = plot_dict["lineage"]["df"]
    lineages = list(lineages_df.columns)
    lineages.remove("epiweek")

    status_counts = {
        status: {"sequences": 0, "lineages": 0} for status in RECOMBINANT_STATUS
    }

    status_seq_df = plot_dict["status"]["df"]
    statuses = list(status_seq_df.columns)
    statuses.remove("epiweek")

    for status in statuses:
        status = status.lower()
        status_lin_df = plot_dict[status]["df"]
        status_lin = list(status_lin_df.columns)
        status_lin.remove("epiweek")

        num_status_lin = len(status_lin)
        status_counts[status]["lineages"] += num_status_lin

        for lin in status_lin:
            num_status_seq = int(sum(lineages_df[lin]))
            status_counts[status]["sequences"] += num_status_seq

    # Number of lineages and sequences
    total_sequences = int(sum([status_counts[s]["sequences"] for s in status_counts]))
    total_lineages = int(sum([status_counts[s]["lineages"] for s in status_counts]))

    # Construct the summary text
    summary = "\n"

    summary += "There are {total_lineages} recombinant lineages".format(
        total_lineages=total_lineages
    )
    # Whether we need a footnote for singletons
    if singletons:
        summary += ".\n"
    else:
        summary += "*.\n"

    for status in RECOMBINANT_STATUS:
        if status in status_counts:
            count = status_counts[status]["lineages"]
        else:
            count = 0
        summary += "  - {lineages} lineages are {status}.\n".format(
            lineages=count, status=status
        )

    summary += "\n"
    summary += "There are {total_sequences} recombinant sequences".format(
        total_sequences=total_sequences
    )

    # Whether we need a footnote for singletons
    if singletons:
        summary += ".\n"
    else:
        summary += "*.\n"

    for status in RECOMBINANT_STATUS:
        if status in status_counts:
            count = status_counts[status]["sequences"]
        else:
            count = 0
        summary += "  - {sequences} sequences are {status}.\n".format(
            sequences=count, status=status
        )
    if not singletons:
        summary += "\n"
        summary += "*Excluding singleton lineages (N=1)"

    body.text_frame.text = summary

    # Adjust font size of body
    for paragraph in body.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = pptx.util.Pt(14)

    # ---------------------------------------------------------------------
    # Status Summary

    for status in RECOMBINANT_STATUS:

        status = status.lower()
        if status not in plot_dict:
            continue

        plot_path = plot_dict[status]["plot_path"]

        # Slide setup
        graph_slide_layout = presentation.slide_layouts[8]
        slide = presentation.slides.add_slide(graph_slide_layout)
        title = slide.shapes.title

        title.text_frame.text = status.title()
        title.text_frame.paragraphs[0].font.bold = True

        chart_placeholder = slide.placeholders[1]
        chart_placeholder.insert_picture(plot_path)
        body = slide.placeholders[2]

        # Stats
        status_df = plot_dict[status]["df"]
        status_lineages = list(status_df.columns)
        status_lineages.remove("epiweek")

        # Order columns
        status_counts = {
            lineage: int(sum(status_df[lineage])) for lineage in status_lineages
        }
        status_lineages = sorted(status_counts, key=status_counts.get, reverse=True)
        num_status_lineages = len(status_lineages)

        summary = "\n"
        summary += "There are {num_status_lineages} {status} lineages.\n".format(
            status=status, num_status_lineages=num_status_lineages
        )

        for lineage in status_lineages:
            seq_count = int(sum(status_df[lineage].dropna()))
            summary += "  - {lineage} ({seq_count})\n".format(
                lineage=lineage, seq_count=seq_count
            )

        body.text_frame.text = summary

        # Adjust font size of body
        for paragraph in body.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = pptx.util.Pt(14)

    # ---------------------------------------------------------------------
    # Geographic Summary

    plot_path = plot_dict["geography"]["plot_path"]
    geo_df = plot_dict["geography"]["df"]
    geos = list(geo_df.columns)
    geos.remove("epiweek")
    # Order columns
    geos_counts = {region: int(sum(geo_df[region])) for region in geos}
    geos = sorted(geos_counts, key=geos_counts.get, reverse=True)

    num_geos = len(geos)

    graph_slide_layout = presentation.slide_layouts[8]
    slide = presentation.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title

    title.text_frame.text = "Geography"
    title.text_frame.paragraphs[0].font.bold = True

    chart_placeholder = slide.placeholders[1]
    chart_placeholder.insert_picture(plot_path)
    body = slide.placeholders[2]

    summary = "\n"
    summary += "Recombinants are observed in {num_geos} {geo}.\n".format(
        num_geos=num_geos, geo=geo
    )

    for region in geos:
        seq_count = int(sum(geo_df[region].dropna()))
        summary += "  - {region} ({seq_count})\n".format(
            region=region, seq_count=seq_count
        )

    body.text_frame.text = summary

    # Adjust font size of body
    for paragraph in body.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = pptx.util.Pt(14)

    # ---------------------------------------------------------------------
    # Largest Summary

    plot_path = plot_dict["largest"]["plot_path"]
    largest_df = plot_dict["largest"]["df"]

    largest_geos = list(largest_df.columns)
    largest_geos.remove("epiweek")

    # Order columns
    largest_geos_counts = {
        region: int(sum(largest_df[region])) for region in largest_geos
    }
    largest_geos = sorted(
        largest_geos_counts, key=largest_geos_counts.get, reverse=True
    )

    num_geos = len(largest_geos)

    largest_lineage = plot_dict["largest"]["lineage"]
    largest_lineage_size = 0

    for lineage in lineages:
        seq_count = int(sum(lineages_df[lineage].dropna()))
        if seq_count > largest_lineage_size:
            largest_lineage_size = seq_count

    graph_slide_layout = presentation.slide_layouts[8]
    slide = presentation.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title

    title.text_frame.text = "Largest"
    title.text_frame.paragraphs[0].font.bold = True

    chart_placeholder = slide.placeholders[1]
    chart_placeholder.insert_picture(plot_path)
    body = slide.placeholders[2]

    summary = "\n"
    summary += "The largest lineage is {lineage} (N={size}).\n".format(
        lineage=largest_lineage,
        size=largest_lineage_size,
    )

    summary += "{lineage} is observed in {num_geo} {geo}.\n".format(
        lineage=largest_lineage, num_geo=num_geos, geo=geo
    )

    for region in largest_geos:
        seq_count = int(sum(largest_df[region].dropna()))
        summary += "  - {region} ({seq_count})\n".format(
            region=region, seq_count=seq_count
        )

    body.text_frame.text = summary

    # Adjust font size of body
    for paragraph in body.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = pptx.util.Pt(14)

    # ---------------------------------------------------------------------
    # RBD Levels

    plot_path = plot_dict["rbd_level"]["plot_path"]
    rbd_df = plot_dict["rbd_level"]["df"]
    rbd_levels = list(rbd_df.columns)
    rbd_levels.remove("epiweek")
    # Order columns
    rbd_counts = {level: int(sum(rbd_df[level])) for level in rbd_levels}
    rbd_levels = dict(sorted(rbd_counts.items()))

    num_rbd_levels = len(rbd_levels)

    graph_slide_layout = presentation.slide_layouts[8]
    slide = presentation.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title

    title.text_frame.text = "Receptor Binding Domain"
    title.text_frame.paragraphs[0].font.bold = True

    chart_placeholder = slide.placeholders[1]
    chart_placeholder.insert_picture(plot_path)
    body = slide.placeholders[2]

    summary = "\n"
    summary += "{num_rbd_levels} RBD levels are observed.\n".format(
        num_rbd_levels=num_rbd_levels,
    )

    for level in rbd_levels:
        seq_count = int(sum(rbd_df[level].dropna()))
        summary += "  - Level {level} ({seq_count})\n".format(
            level=level, seq_count=seq_count
        )

    body.text_frame.text = summary

    # Adjust font size of body
    for paragraph in body.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = pptx.util.Pt(14)

    # ---------------------------------------------------------------------
    # Parents Summary (Clade)

    plot_path = plot_dict["parents_clade"]["plot_path"]
    parents_df = plot_dict["parents_clade"]["df"]

    parents = list(parents_df.columns)
    parents.remove("epiweek")

    parents_counts = {p: int(sum(parents_df[p])) for p in parents}
    parents = sorted(parents_counts, key=parents_counts.get, reverse=True)

    num_parents = len(parents)

    graph_slide_layout = presentation.slide_layouts[8]
    slide = presentation.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title

    title.text_frame.text = "Parents (Clade)"
    title.text_frame.paragraphs[0].font.bold = True

    chart_placeholder = slide.placeholders[1]
    chart_placeholder.insert_picture(plot_path)
    body = slide.placeholders[2]

    summary = "\n"
    summary += "There are {num_parents} clade combinations.\n".format(
        num_parents=num_parents
    )

    for parent in parents:
        seq_count = int(sum(parents_df[parent].dropna()))
        summary += "  - {parent} ({seq_count})\n".format(
            parent=parent, seq_count=seq_count
        )

    body.text_frame.text = summary

    # Adjust font size of body
    for paragraph in body.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = pptx.util.Pt(14)

    # ---------------------------------------------------------------------
    # Parents Summary (Lineage)

    plot_path = plot_dict["parents_lineage"]["plot_path"]
    parents_df = plot_dict["parents_lineage"]["df"]

    parents = list(parents_df.columns)
    parents.remove("epiweek")

    parents_counts = {p: int(sum(parents_df[p])) for p in parents}
    parents = sorted(parents_counts, key=parents_counts.get, reverse=True)

    num_parents = len(parents)

    graph_slide_layout = presentation.slide_layouts[8]
    slide = presentation.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title

    title.text_frame.text = "Parents (Lineage)"
    title.text_frame.paragraphs[0].font.bold = True

    chart_placeholder = slide.placeholders[1]
    chart_placeholder.insert_picture(plot_path)
    body = slide.placeholders[2]

    summary = "\n"
    summary += "There are {num_parents} lineage combinations.\n".format(
        num_parents=num_parents
    )

    for parent in parents:
        seq_count = int(sum(parents_df[parent].dropna()))
        summary += "  - {parent} ({seq_count})\n".format(
            parent=parent, seq_count=seq_count
        )

    body.text_frame.text = summary

    # Adjust font size of body
    for paragraph in body.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = pptx.util.Pt(14)

    # ---------------------------------------------------------------------
    # Breakpoints Clade Summary

    plot_path = plot_dict["breakpoints_clade"]["plot_path"]

    graph_slide_layout = presentation.slide_layouts[8]
    slide = presentation.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title

    title.text_frame.text = "Breakpoints (Clade)"
    title.text_frame.paragraphs[0].font.bold = True

    chart_placeholder = slide.placeholders[1]
    chart_placeholder.insert_picture(plot_path)
    body = slide.placeholders[2]

    # ---------------------------------------------------------------------
    # Breakpoints Lineage Summary

    plot_path = plot_dict["breakpoints_lineage"]["plot_path"]

    graph_slide_layout = presentation.slide_layouts[8]
    slide = presentation.slides.add_slide(graph_slide_layout)
    title = slide.shapes.title

    title.text_frame.text = "Breakpoints (Lineage)"
    title.text_frame.paragraphs[0].font.bold = True

    chart_placeholder = slide.placeholders[1]
    chart_placeholder.insert_picture(plot_path)
    body = slide.placeholders[2]

    # ---------------------------------------------------------------------
    # Footer

    # Versions
    pipeline_version = linelist_df["pipeline_version"].values[0]
    recombinant_classifier_dataset = linelist_df[
        "recombinant_classifier_dataset"
    ].values[0]

    for slide in presentation.slides:

        # Don't add footer to title slide
        if slide.slide_layout == presentation.slide_layouts[0]:
            continue

        footer = slide.shapes.add_shape(
            autoshape_type_id=MSO_SHAPE.RECTANGLE,
            left=pptx.util.Cm(3),
            top=pptx.util.Cm(17),
            width=pptx.util.Pt(800),
            height=pptx.util.Pt(30),
        )

        footer.fill.solid()
        footer.fill.fore_color.rgb = pptx.dml.color.RGBColor(255, 255, 255)
        footer.line.color.rgb = pptx.dml.color.RGBColor(0, 0, 0)
        p = footer.text_frame.paragraphs[0]
        p.text = "{} {} {}".format(
            pipeline_version, " " * 20, recombinant_classifier_dataset
        )
        p.font.size = pptx.util.Pt(14)
        p.font.color.rgb = pptx.dml.color.RGBColor(0, 0, 0)

    # ---------------------------------------------------------------------
    # Saving file
    presentation.save(output)


if __name__ == "__main__":
    main()
